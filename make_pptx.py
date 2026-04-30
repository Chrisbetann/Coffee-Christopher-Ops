from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

BROWN  = RGBColor(0x6B, 0x3F, 0x1E)
TAN    = RGBColor(0xC9, 0xA9, 0x6E)
CREAM  = RGBColor(0xF5, 0xEF, 0xE6)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
DARK   = RGBColor(0x2C, 0x1A, 0x0E)
GRAY   = RGBColor(0x6B, 0x72, 0x80)
LGRAY  = RGBColor(0xF3, 0xF4, 0xF6)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def box(slide, l, t, w, h, color, radius=False):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def txt(slide, text, l, t, w, h, size, color, bold=False, align=PP_ALIGN.LEFT, wrap=True):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    return tb

def bullet_slide(prs, title, subtitle, bullets, icon=""):
    slide = blank(prs)
    bg(slide, CREAM)
    # Left accent bar
    box(slide, 0, 0, 0.08, 7.5, BROWN)
    # Header area
    box(slide, 0.08, 0, 13.25, 1.6, BROWN)
    txt(slide, icon + "  " + title if icon else title, 0.4, 0.18, 12, 0.8, 32, WHITE, bold=True)
    txt(slide, subtitle, 0.4, 0.95, 12, 0.55, 13, TAN)
    # Bullet cards
    y = 1.85
    for i, (head, body) in enumerate(bullets):
        card = box(slide, 0.4, y, 12.5, 0.78, WHITE)
        # Dot
        dot = box(slide, 0.55, y + 0.25, 0.18, 0.18, TAN)
        txt(slide, head, 0.9, y + 0.04, 11.7, 0.35, 13, DARK, bold=True)
        txt(slide, body, 0.9, y + 0.38, 11.7, 0.35, 11, GRAY)
        y += 0.88
    return slide

def two_col(prs, title, subtitle, left_items, right_items, left_head, right_head):
    slide = blank(prs)
    bg(slide, CREAM)
    box(slide, 0, 0, 0.08, 7.5, BROWN)
    box(slide, 0.08, 0, 13.25, 1.6, BROWN)
    txt(slide, title, 0.4, 0.18, 12, 0.8, 32, WHITE, bold=True)
    txt(slide, subtitle, 0.4, 0.95, 12, 0.55, 13, TAN)
    # Two columns
    for col, (head, items) in enumerate([(left_head, left_items), (right_head, right_items)]):
        x = 0.4 + col * 6.5
        box(slide, x, 1.75, 6.1, 0.42, BROWN)
        txt(slide, head, x + 0.15, 1.8, 5.8, 0.35, 13, WHITE, bold=True)
        y = 2.3
        for item in items:
            box(slide, x, y, 6.1, 0.65, WHITE)
            txt(slide, "▸  " + item, x + 0.15, y + 0.08, 5.8, 0.5, 12, DARK)
            y += 0.72
    return slide

# ── SLIDE 1: Title ──────────────────────────────────────────────────────────
slide = blank(prs)
bg(slide, BROWN)
box(slide, 0, 5.5, 13.33, 2.0, DARK)
box(slide, 0, 0, 0.35, 7.5, TAN)
txt(slide, "☕", 0.6, 0.5, 3, 2.0, 80, WHITE)
txt(slide, "Coffee Christopher", 3.0, 0.9, 10, 1.1, 44, WHITE, bold=True)
txt(slide, "Operations Suite", 3.0, 1.95, 10, 0.8, 30, TAN)
txt(slide, "Full-Stack Ops Platform  ·  Loyalty Program  ·  PWA  ·  Admin Portal", 3.0, 2.85, 10, 0.6, 14, RGBColor(0xFF,0xFF,0xFF))
txt(slide, "Sprint Planning Review  ·  April 2026", 0.5, 6.1, 12, 0.6, 13, TAN)
txt(slide, "Built with React · Node.js · Prisma · PostgreSQL (Supabase) · Tailwind CSS", 0.5, 6.65, 12, 0.5, 11, GRAY)

# ── SLIDE 2: Project Overview ───────────────────────────────────────────────
bullet_slide(prs,
    "Project Overview", "What is Coffee Christopher Ops?",
    [
        ("What it is", "A full-stack operations platform for a real coffee shop — menu, orders, loyalty, inventory, analytics, and promotions."),
        ("Who uses it", "Two audiences: customers (public menu, loyalty card, order tracking) and admin staff (full back-office portal)."),
        ("Why we built it", "Replace manual pen-and-paper loyalty tracking and spreadsheet inventory with a single, mobile-friendly web app."),
        ("Deployment model", "Vite + React frontend served via Node.js/Express backend, PostgreSQL database hosted on Supabase, deployable as a PWA."),
        ("Current status", "Core features fully operational end-to-end: menu loads, orders placed, loyalty enrolled, promos sent, reminders emailed."),
    ],
    icon="📋"
)

# ── SLIDE 3: Tech Stack ─────────────────────────────────────────────────────
two_col(prs,
    "Tech Stack", "Frontend + Backend + Infrastructure",
    ["React 18 + Vite 5", "React Router v6", "Tailwind CSS (brand tokens)", "Axios for API calls", "QR Code rendering (qrcode.react)", "PWA: Service Worker + Web Manifest"],
    ["Node.js + Express", "Prisma ORM (PostgreSQL)", "Supabase (hosted Postgres)", "JWT authentication (jsonwebtoken)", "Nodemailer (SMTP via Gmail)", "bcryptjs for password hashing"],
    "Frontend", "Backend"
)

# ── SLIDE 4: Feature Map ────────────────────────────────────────────────────
bullet_slide(prs,
    "Features Built", "Eight core modules across customer and admin surfaces",
    [
        ("☕  Menu", "Public browsable menu with categories, item details, available/unavailable toggle, and photo support."),
        ("🛒  Orders", "Customer order placement with confirmation screen; admin order queue with status updates (pending → in-progress → complete)."),
        ("🎟  Loyalty Program v2", "Phone-number lookup, QR code stamp card, registration, admin stamp/redeem, customer wallet-ready card."),
        ("📣  Promotions", "Admin creates promos (% off, $ off, item discount), views recipient list, sends individual emails via SMTP."),
        ("📦  Inventory", "Ingredient tracking with low-stock alerts, audit log per ingredient, admin CRUD."),
        ("⭐  Reviews", "Post-order review submission, per-item ratings, admin moderation and delete."),
        ("📊  Dashboard", "Sales totals, top items, order volume chart — filterable by day/week/month."),
    ],
    icon="🗺"
)

# ── SLIDE 5: Loyalty Program Deep Dive ─────────────────────────────────────
bullet_slide(prs,
    "Loyalty Program v2", "The centrepiece feature — every 6th drink free",
    [
        ("Customer flow", "Visit /loyalty → enter phone number → view stamp card with QR code → share or install as PWA to home screen."),
        ("Stamp & redeem", "Admin scans QR code → adds stamp or redeems free drink. Stamps roll over after redemption."),
        ("Registration", "New customers register at /loyalty/register — first name, last name, email, phone — card created instantly."),
        ("Weekly reminders", "Admin clicks one button → server sends personalised emails to all members, randomly picking from 4 different message templates."),
        ("PWA install prompt", "iOS share-sheet and Android install banner appear on the loyalty card page so customers can add the card to their home screen."),
    ],
    icon="🎟"
)

# ── SLIDE 6: Promotions Module ──────────────────────────────────────────────
bullet_slide(prs,
    "Promotions Module", "Create, track, and send targeted offers",
    [
        ("Creating a promo", "Admin fills title, optional description, discount type (% off / $ off / $ off item) and value — stored in database instantly."),
        ("Recipient list", "Every active loyalty member appears as a recipient automatically — admin sees name, email, phone, and sent status."),
        ("Email sending", "Clicking 📧 Email sends a personalised offer email via Gmail SMTP directly from the server — no email client needed."),
        ("SMS", "Clicking 📱 SMS opens the Messages app with a pre-filled text message ready to send from the admin's phone."),
        ("Send tracking", "Each email/SMS send is tracked per-customer per-promo with a timestamp. Checkbox auto-ticks after server confirms delivery."),
    ],
    icon="📣"
)

# ── SLIDE 7: PWA & Service Worker ──────────────────────────────────────────
bullet_slide(prs,
    "PWA & Offline Support", "Installable, offline-capable loyalty card",
    [
        ("Web App Manifest", "manifest.webmanifest with brand colours, maskable icons (192px + 512px SVG), and display: standalone for full-screen install."),
        ("Service Worker", "Cache-first shell (/, /loyalty, /menu), network-first API calls with loyalty card fallback, stale-while-revalidate for static assets."),
        ("Install prompt", "Custom InstallPrompt component detects beforeinstallprompt (Android) and shows iOS share-sheet instructions — dismissible for 7 days."),
        ("Dev-mode safety", "Service worker only registers in production (PROD guard in index.jsx) so Vite HMR is never intercepted during development."),
        ("Share button", "Native Web Share API with clipboard fallback — lets customers share their loyalty card URL directly from the card page."),
    ],
    icon="📱"
)

# ── SLIDE 8: Key Challenges Solved ─────────────────────────────────────────
bullet_slide(prs,
    "Key Challenges Solved", "What we debugged and fixed during this sprint",
    [
        ("Dual repo confusion", "Two clones of the repo (~/Coffee-Christopher-Ops vs ~/Desktop/...) caused the wrong server to run — diagnosed via file path in Prisma errors."),
        ("Supabase auth failures", "Database password rotation without updating both .env files caused P1001/auth errors — fixed by sed-replacing both copies."),
        ("Service worker crashing Vite", "SW fetch handler returned undefined, breaking Vite HMR — fixed with PROD guard and Response.error() fallbacks."),
        ("mailto: links doing nothing", "macOS had no default email client; replaced all mailto: hrefs with server-side nodemailer SMTP sends through Gmail."),
        ("Stale browser cache", "Old JS bundle cached by Vite and Opera after code updates — fixed by killing PID, clearing node_modules/.vite, and switching to Chrome."),
    ],
    icon="🔧"
)

# ── SLIDE 9: Email Infrastructure ──────────────────────────────────────────
bullet_slide(prs,
    "Email Infrastructure", "Server-side SMTP sending with Nodemailer",
    [
        ("Transport", "Nodemailer configured with Gmail SMTP (smtp.gmail.com:587, STARTTLS) using a Gmail App Password — not the account password."),
        ("Promo emails", "POST /api/promos/:id/recipients/:custId/email — personalises subject + body, sends, then auto-marks email_sent in database."),
        ("Weekly reminders", "POST /api/loyalty/admin/reminders/weekly — sends to ALL loyalty members, 4 random templates, skips stamp template for 0-stamp members."),
        ("4 reminder templates", '"We miss you" · "Time for a great coffee" · "Your stamps are building up" · "Come say hi" — picked randomly per customer.'),
        ("Error handling", "Promise.allSettled ensures one bad email address doesn't stop the rest — admin sees exact count of sent vs failed."),
    ],
    icon="📧"
)

# ── SLIDE 10: Database Schema ───────────────────────────────────────────────
two_col(prs,
    "Database Schema", "PostgreSQL via Prisma ORM on Supabase",
    ["Category", "MenuItem", "Order + OrderItem", "Admin (JWT auth)", "Review"],
    ["Customer (loyalty)", "PromoSend (tracking)", "Promo", "Ingredient (inventory)", "IngredientLog (audit)"],
    "Core Tables", "Loyalty & Ops Tables"
)

# ── SLIDE 11: Roadmap ───────────────────────────────────────────────────────
bullet_slide(prs,
    "Roadmap & Next Steps", "What comes after this sprint",
    [
        ("Password security", "Rotate Supabase DB password and Gmail App Password exposed during debugging — store secrets in a proper secrets manager."),
        ("Consolidate repos", "Delete ~/Coffee-Christopher-Ops (home folder clone) — work exclusively from ~/Desktop/Coffee-Christopher-Ops to avoid confusion."),
        ("Production deployment", "Deploy frontend to Vercel/Netlify, backend to Railway/Render — set environment variables in hosting dashboard instead of .env files."),
        ("Apple & Google Wallet", "Wallet pass endpoints exist but need APPLE_PASS_CERT and GOOGLE_WALLET_SERVICE_ACCOUNT credentials to activate."),
        ("Automated scheduling", "Replace manual 'Send Reminder' button with a cron job (node-cron or Supabase Edge Functions) to send reminders automatically weekly."),
    ],
    icon="🚀"
)

# ── SLIDE 12: Thank You ─────────────────────────────────────────────────────
slide = blank(prs)
bg(slide, BROWN)
box(slide, 0, 5.8, 13.33, 1.7, DARK)
box(slide, 0, 0, 0.35, 7.5, TAN)
txt(slide, "☕", 0.6, 1.2, 2, 1.5, 72, WHITE)
txt(slide, "Coffee Christopher Ops", 3.0, 1.2, 10, 0.9, 38, WHITE, bold=True)
txt(slide, "Built for real. Built to last.", 3.0, 2.1, 10, 0.6, 20, TAN)
txt(slide, "React  ·  Node.js  ·  Prisma  ·  Supabase  ·  Tailwind  ·  Nodemailer  ·  PWA", 3.0, 2.85, 10, 0.5, 13, RGBColor(0xFF,0xFF,0xFF))
txt(slide, "github.com/Chrisbetann/Coffee-Christopher-Ops", 3.0, 3.5, 10, 0.5, 13, TAN)
txt(slide, "Sprint Planning Review  ·  April 2026", 0.5, 6.1, 12, 0.5, 13, TAN)

out = "/home/user/Coffee-Christopher-Ops/Coffee-Christopher-Ops-Presentation.pptx"
prs.save(out)
print(f"Saved: {out}")
